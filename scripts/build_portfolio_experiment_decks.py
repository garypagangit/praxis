from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


OUT_DIR = Path("paper/portfolio_experiment_decks")


COLORS = {
    "ink": RGBColor(17, 24, 39),
    "muted": RGBColor(75, 85, 99),
    "blue": RGBColor(37, 99, 235),
    "blue_soft": RGBColor(219, 234, 254),
    "green_soft": RGBColor(220, 252, 231),
    "red_soft": RGBColor(254, 226, 226),
    "gray_soft": RGBColor(249, 250, 251),
}


def add_textbox(slide, left, top, width, height, text, size=22, bold=False, color="ink"):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = COLORS[color]
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.55), Inches(0.32), Inches(12.25), Inches(0.6), title, 28, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.02), Inches(12.15), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["blue"]
    line.line.color.rgb = COLORS["blue"]
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(1.15), Inches(11.8), Inches(0.42), subtitle, 14, False, "muted")


def add_bullets(slide, items, left=0.75, top=1.75, width=8.6, height=4.8, size=21):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS["ink"]


def add_table(slide, rows, left, top, width, height, font_size=13):
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height)).table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.bold = r == 0
                p.font.color.rgb = COLORS["ink"]
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["blue_soft"]
    return table


def add_card(slide, left, top, width, height, title, body, fill="gray_soft"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.color.rgb = RGBColor(209, 213, 219)
    add_textbox(slide, Inches(left + 0.18), Inches(top + 0.15), Inches(width - 0.36), Inches(0.32), title, 15, True)
    add_textbox(slide, Inches(left + 0.18), Inches(top + 0.55), Inches(width - 0.36), Inches(height - 0.7), body, 13)


def add_note(slide, note):
    slide.notes_slide.notes_text_frame.text = note


def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def slide(prs, title, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, title, subtitle)
    return s


def save(prs, path):
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    print(path)


def build_attack_ttp_retrieval():
    prs = new_deck()
    path = OUT_DIR / "EXPERIMENT_03_ATTACK_TTP_RETRIEVAL_DECK_20260514.pptx"

    s = slide(prs, "Experiment 3: ATT&CK TTP-Set Profile Retrieval", "Second narrow positive candidate")
    add_bullets(s, [
        "Goal: rank likely ATT&CK group profiles from a small observed set of techniques.",
        "Current status: selected as result #2, with scope narrowed away from prose attribution.",
        "Headline: 5-shot top-5 accuracy is 0.960 for overlap and 0.879 for SVD.",
    ], width=11.5, size=24)
    add_note(s, "Lead with scope discipline: this is profile retrieval, not actor authorship.")

    s = slide(prs, "Paper Anchor And Gap")
    add_card(s, 0.75, 1.75, 5.8, 2.1, "Primary anchor", "Strom et al. (2020), MITRE ATT&CK: Design and philosophy. ATT&CK provides empirically grounded adversary technique profiles.", "blue_soft")
    add_card(s, 6.85, 1.75, 5.8, 2.1, "Gap", "There is no formal few-shot retrieval protocol that asks how many observed TTPs are enough to recover likely group profiles.", "green_soft")
    add_card(s, 0.75, 4.2, 11.9, 1.35, "Boundary", "Hamilton et al. (2017) motivates graph embeddings, but the local GraphSAGE pilot failed. The selected result is the simple retrieval protocol.", "gray_soft")

    s = slide(prs, "Protocol")
    add_table(s, [
        ["Step", "Decision"],
        ["Data", "ATT&CK group-technique matrix: 174 groups, 697 techniques, 4546 edges"],
        ["Eligible groups", "121 groups with degree >= 10"],
        ["Queries", "Sample k techniques from a known group, k in {1, 3, 5, 10}"],
        ["Metrics", "Top-1, Top-5, Top-10, MRR, median rank"],
        ["Scope", "Profile retrieval from TTP sets, not CTI prose attribution"],
    ], 0.75, 1.55, 11.9, 4.7, 15)

    s = slide(prs, "Main Result With Floor Baselines")
    add_table(s, [
        ["Method", "Shots", "Top-1", "Top-5", "MRR", "Median rank", "Lift vs random"],
        ["random_uniform", "5", "0.005", "0.028", "0.033", "84.0", "1.0x"],
        ["frequency_prior", "5", "0.008", "0.041", "0.044", "61.0", "1.5x"],
        ["overlap_cosine", "5", "0.721", "0.960", "0.824", "1.0", "34.2x"],
        ["svd32_graph_embedding", "5", "0.623", "0.879", "0.732", "1.0", "31.3x"],
    ], 0.75, 1.55, 11.9, 4.5, 15)
    add_note(s, "The closeout claim is now stronger: retrieval beats random and frequency-prior floors while preserving the original overlap/SVD numbers.")

    s = slide(prs, "Degree-Bucket Check")
    add_table(s, [
        ["Method", "Bucket", "Top-5", "MRR", "Median rank"],
        ["overlap_cosine", "low", "0.995", "0.953", "1.0"],
        ["overlap_cosine", "mid", "0.995", "0.866", "1.0"],
        ["overlap_cosine", "high", "0.887", "0.644", "2.0"],
        ["svd32_graph_embedding", "low", "0.951", "0.866", "1.0"],
        ["svd32_graph_embedding", "mid", "0.854", "0.667", "1.0"],
        ["svd32_graph_embedding", "high", "0.831", "0.659", "1.0"],
    ], 0.75, 1.55, 11.9, 4.5, 15)

    s = slide(prs, "Negative GNN Gate")
    add_table(s, [
        ["Mode", "Method", "Shots", "Top-5", "Median rank"],
        ["known_profile", "overlap_cosine_train", "5", "0.985", "1.0"],
        ["known_profile", "svd32_train_graph", "5", "0.926", "1.0"],
        ["known_profile", "graphsage_linkpred", "5", "0.060", "60.0"],
        ["held_edge", "graphsage_linkpred", "5", "0.073", "34.0"],
    ], 0.75, 1.55, 11.9, 4.5, 15)

    s = slide(prs, "Allowed Claims")
    add_bullets(s, [
        "Allowed: five observed ATT&CK techniques can retrieve correct group profiles with high top-5 accuracy.",
        "Allowed: SVD group-technique embeddings preserve useful retrieval signal.",
        "Not allowed: GraphSAGE is better than simple baselines.",
        "Not allowed: CTI prose attribution or actor authorship is solved.",
    ], width=11.5, size=23)

    s = slide(prs, "Next Work")
    add_bullets(s, [
        "Add random and frequency-prior baselines.",
        "Add degree-bucket analysis to test whether high-degree groups dominate performance.",
        "Add example retrievals: query TTPs, top candidates, and why the correct group ranked where it did.",
        "Convert the result report into a compact thesis section.",
    ], width=11.5, size=23)

    save(prs, path)


def build_provenance_labels():
    prs = new_deck()
    path = OUT_DIR / "EXPERIMENT_04_PROVENANCE_LABEL_PATH_DECK_20260514.pptx"

    s = slide(prs, "Experiment 4: Provenance Labels / OpTC Path", "Architecture is ready; labels are the blocker")
    add_bullets(s, [
        "Goal: reopen provenance drift, TGN, SSL, graph routing, MIA, and watermarking with honest interval labels.",
        "Current status: architecture ready but supervised claims are blocked.",
        "Decision: targeted OpTC subset first; Cadets interval labels as fallback.",
    ], width=11.5, size=24)

    s = slide(prs, "Paper Anchor And Gap")
    add_card(s, 0.75, 1.65, 5.8, 2.0, "Primary anchor", "Han et al. (2020), UNICORN. Provenance can support runtime APT detection.", "blue_soft")
    add_card(s, 6.85, 1.65, 5.8, 2.0, "Drift anchor", "Gama et al. (2014), concept drift adaptation. Chronological detector stability needs labels.", "blue_soft")
    add_card(s, 0.75, 4.0, 11.9, 1.65, "Gap", "Most of the current portfolio is blocked not by graph modeling, but by missing benign/attack interval truth for provenance windows.", "green_soft")

    s = slide(prs, "Why Cadets Is Blocked")
    add_table(s, [
        ["Source", "Windows", "Attack-touch", "Benign / unlabeled", "Decision"],
        ["Full E5 Cadets", "9611", "9609", "2", "No supervised detector claim"],
    ], 0.75, 1.65, 11.9, 1.3, 15)
    add_bullets(s, [
        "Node-touch labels are too broad: nearly every window touches attack-labeled nodes.",
        "Density proxy is learnable but not ground truth.",
        "Training a supervised detector here would create a bogus claim.",
    ], top=3.35, width=11.5, size=22)

    s = slide(prs, "Minimum Label Gate")
    add_table(s, [
        ["Check", "Minimum"],
        ["Real classes", "At least 2"],
        ["Benign windows", ">= 20"],
        ["Attack/anomaly windows", ">= 20"],
        ["Stage-specific claim", ">= 20 windows for that stage"],
        ["Split support", "Train/validation/test contain claimed positive class"],
        ["Label source", "Interval truth, confirmed benign spans, or trusted release annotation"],
    ], 0.75, 1.55, 11.9, 4.6, 15)

    s = slide(prs, "Candidate Paths")
    add_table(s, [
        ["Path", "Strength", "Risk", "Decision"],
        ["Cadets intervals", "Reuses 480M-event run", "Manual labels may be weak", "Fallback"],
        ["Another stream", "Could avoid label collapse", "Discovery/setup unknown", "Investigate later"],
        ["Targeted OpTC subset", "Red-team ground truth plus benign activity", "Schema/setup cost", "Recommended first"],
    ], 0.75, 1.55, 11.9, 4.25, 14)

    s = slide(prs, "OpTC Feasibility Check")
    add_table(s, [
        ["Artifact", "Result"],
        ["Ground-truth parser", "Extracts timestamped seed events from OpTC red-team PDF"],
        ["Timestamped red-team events", "101"],
        ["Covered days", "Plain PowerShell Empire; Custom Powershell Empire; Malicious Upgrade"],
        ["Next blocker", "Attach eCAR host telemetry around seed timestamps"],
    ], 0.75, 1.55, 11.9, 4.3, 15)

    s = slide(prs, "What Opens After Gate")
    add_bullets(s, [
        "Concept drift: chronological detector stability with real labels.",
        "TGN/SSL: representation experiments with meaningful positive and negative windows.",
        "Watermarking/MIA/adversarial robustness: only after stable detector families exist.",
        "Stage routing: only if stage intervals or reliable stage labels exist.",
    ], width=11.5, size=23)

    s = slide(prs, "Next Action")
    add_bullets(s, [
        "Download or mirror one targeted OpTC eCAR shard, not the whole release.",
        "Run scripts/build_optc_window_gate.ps1.",
        "Stop unless the gate finds >= 20 benign and >= 20 attack/anomaly windows.",
        "If the gate passes, train detector-zoo baselines; if it fails, archive the blocker.",
    ], width=11.5, size=23)

    save(prs, path)


def build_sec_lord():
    prs = new_deck()
    path = OUT_DIR / "EXPERIMENT_07_SEC_LORD_REDESIGN_GATE_DECK_20260514.pptx"

    s = slide(prs, "Experiment 7: SEC-LoRD / DS-LoRD", "Current method is negative; redesign gate only")
    add_bullets(s, [
        "Goal: test whether security-domain evidence improves CTI task behavior before any extraction claim.",
        "Current broad domain-seeded prompting failed under strict parsing.",
        "Decision: no extraction until a retrieved-evidence prompt clears a strict gate.",
    ], width=11.5, size=24)

    s = slide(prs, "Paper Anchor And Gap")
    add_card(s, 0.75, 1.65, 5.8, 2.0, "Extraction anchor", "Carlini et al. (2021). LLMs can leak or reproduce training data under extraction-style prompting.", "blue_soft")
    add_card(s, 6.85, 1.65, 5.8, 2.0, "RAG anchor", "Lewis et al. (2020). Retrieval-augmented generation separates evidence from parametric recall.", "blue_soft")
    add_card(s, 0.75, 4.0, 11.9, 1.65, "Gap", "Does question-specific CTI evidence improve strict answer compliance before any LoRD-style extraction experiment is justified?", "green_soft")

    s = slide(prs, "Current Strict Audit")
    add_table(s, [
        ["Model", "Vanilla strict acc", "Seeded strict acc", "Delta"],
        ["Llama-3.2-3B-Instruct", "0.276", "0.090", "-0.186"],
        ["Llama-3.1-8B-Instruct", "0.466", "0.284", "-0.182"],
    ], 0.75, 1.65, 11.9, 2.1, 16)
    add_bullets(s, [
        "Broad seed stuffing reduced task compliance.",
        "Seeded prompts increased invalid answers.",
        "The current method is stopped, not rescued.",
    ], top=4.1, width=11.5, size=22)

    s = slide(prs, "Why The Original Method Failed")
    add_bullets(s, [
        "The prompt injected broad domain vocabulary instead of question-specific evidence.",
        "The model could answer in non-strict prose that the old parser would over-credit.",
        "The corrected strict parser exposed the real failure.",
        "Extraction would be meaningless if the attack prompt first harms basic task accuracy.",
    ], width=11.5, size=23)

    s = slide(prs, "New Method")
    add_table(s, [
        ["Component", "Design"],
        ["Evidence retrieval", "Retrieve 1-3 short ATT&CK/CTI facts tied to the question"],
        ["Prompt separation", "Put facts in an Evidence block, not broad domain seed lists"],
        ["Answer format", "Force: Answer: <A|B|C|D>"],
        ["Comparisons", "Vanilla strict, retrieved-evidence prompt, broad seed negative control"],
    ], 0.75, 1.55, 11.9, 4.6, 15)

    s = slide(prs, "Gate Metrics")
    add_table(s, [
        ["Metric", "Threshold"],
        ["Retrieved-evidence strict accuracy vs vanilla", ">= +0.030 absolute"],
        ["Invalid response rate", "<= vanilla invalid rate"],
        ["Seeded negative control", "Reported, not hidden"],
        ["Paired wins", "Evidence-only wins exceed vanilla-only wins"],
        ["Budget", "Local or one small GPU/API batch; max 500 questions per model"],
    ], 0.75, 1.55, 11.9, 4.6, 15)

    s = slide(prs, "Pass / Fail Decision")
    add_card(s, 0.75, 1.75, 5.8, 2.4, "If pass", "Design a separate extraction experiment. The extraction objective must not depend on broad prompt seeding improving CTI accuracy.", "green_soft")
    add_card(s, 6.85, 1.75, 5.8, 2.4, "If fail", "Archive SEC-LoRD as negative for CTI-MCQ and pivot to better-labeled TTP linking or retrieval evaluation.", "red_soft")
    add_card(s, 0.75, 4.55, 11.9, 1.0, "Guardrail", "No model extraction claim until the strict answer-format gate passes.", "gray_soft")

    s = slide(prs, "Next Action")
    add_bullets(s, [
        "Implement retrieved-evidence prompt construction over the same CTI-MCQ set.",
        "Run vanilla vs retrieved-evidence vs broad-seeded negative control.",
        "Score with strict parser only.",
        "Escalate to extraction only if the gate passes.",
    ], width=11.5, size=23)

    save(prs, path)


def main() -> None:
    build_attack_ttp_retrieval()
    build_provenance_labels()
    build_sec_lord()


if __name__ == "__main__":
    main()
