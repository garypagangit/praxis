from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


OUT = Path("paper/praxis06_tta/defense_slides/PRAXIS06_DEFENSE_SLIDES_20260514.pptx")
METHOD_FIG = Path(
    "reports/tta_streaming_apt/paper_assets_20260509/figures/figure1_method_diagram.png"
)
RECON_FIG = Path(
    "reports/tta_streaming_apt/paper_assets_20260509/figures/figure4_recon_pr_operating_points.png"
)
DE_FIG = Path(
    "reports/tta_streaming_apt/paper_assets_20260509/figures/figure5_de_pr_operating_points.png"
)


def add_textbox(slide, left, top, width, height, text, size=22, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(31, 41, 55)
    return box


def add_bullets(slide, items, left=0.75, top=1.65, width=8.8, height=4.5, size=21):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(31, 41, 55)


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.55), Inches(0.35), Inches(12.25), Inches(0.55), title, 30, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.05), Inches(12.15), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(37, 99, 235)
    line.line.color.rgb = RGBColor(37, 99, 235)
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(1.15), Inches(11.8), Inches(0.35), subtitle, 14)


def add_table(slide, rows, left, top, width, height, font_size=13):
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height)).table
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.bold = r == 0
                p.font.color.rgb = RGBColor(17, 24, 39)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(219, 234, 254)
    return table


def add_speaker_note(slide, note):
    notes = slide.notes_slide.notes_text_frame
    notes.text = note


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slides = []

    slide = prs.slides.add_slide(blank)
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(12), Inches(0.8), "Selective Test-Time Adaptation for Streaming APT Stage Detection", 34, True)
    add_textbox(slide, Inches(0.7), Inches(2.35), Inches(12), Inches(0.55), "Under Source-File Shift", 28, True)
    add_textbox(slide, Inches(0.72), Inches(3.25), Inches(10.5), Inches(0.6), "Praxis 06 defense result: narrow, locked, cloud-audited", 20)
    add_textbox(slide, Inches(0.72), Inches(5.7), Inches(10.5), Inches(0.4), "Core claim: safety-gated no-label TTA recovers shifted Reconnaissance without a broad TTA generality claim.", 15)
    add_speaker_note(slide, "Open by saying this is the lead positive result after portfolio triage. The claim is deliberately narrow and evidence-bounded.")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "The Problem", "APT stage detection under deployment-style shift")
    add_bullets(slide, [
        "Source-file, day, host, and traffic shift are deployment realities.",
        "Rare stages can collapse even when the overall detector appears usable.",
        "Reconnaissance is operationally important but fragile under shift.",
        "Security ML needs improvement without damaging high-consequence stages.",
    ])
    add_speaker_note(slide, "Frame this as deployment realism, not benchmark polishing.")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Why Previous Paths Were Not Enough")
    add_bullets(slide, [
        "Praxis 04 stage routing failed under shift.",
        "Treatment-Stage Macro-F1 0.5981 vs Baseline-TSE 0.6313.",
        "Stage-conditioned imbalance did not rescue rare classes cleanly.",
        "Provenance graph detector claims remain label-blocked.",
    ])
    add_speaker_note(slide, "This sets up why we moved from more routing or weighting to deployment-time adaptation.")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Research Questions And Decision Gates")
    add_table(slide, [
        ["Question", "Gate"],
        ["Improve Macro-F1?", "delta >= +0.05"],
        ["Recover Recon?", "delta >= +0.25"],
        ["Preserve DE?", "mean delta >= 0; per-seed >= -0.05"],
        ["Stay selective?", "override <= 0.05 in locked replay"],
        ["More than filtering?", "confidence-reject Recon remains collapsed"],
    ], 0.75, 1.55, 11.8, 3.7, 15)
    add_speaker_note(slide, "Emphasize that gates prevent a broad post-hoc tuning story.")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Method: Selective TTA")
    if METHOD_FIG.exists():
        slide.shapes.add_picture(str(METHOD_FIG), Inches(0.8), Inches(1.55), width=Inches(7.1))
    add_bullets(slide, [
        "Frozen support-floor MLP produces baseline probabilities.",
        "BatchNorm adaptation updates unlabeled stream statistics.",
        "recon_guarded gate accepts only selected adapted decisions.",
        "Confident Data Exfiltration predictions are protected.",
    ], left=8.1, top=1.55, width=4.6, height=4.4, size=17)
    add_speaker_note(slide, "The method is not unconstrained TTA. The safety gate is the contribution.")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Data And Leakage Controls")
    add_table(slide, [
        ["Split", "Rows", "Sources", "Recon", "DE"],
        ["Train", "307,733", "142", "5,151", "3,077"],
        ["Validation", "61,886", "6", "23,791", "2,253"],
        ["Test", "65,869", "25", "5,852", "2,192"],
    ], 0.9, 1.55, 8.7, 2.1, 15)
    add_bullets(slide, [
        "Held-out source-file split.",
        "Zero source/group overlap across splits.",
        "Temporal deltas use reset_each_split.",
        "Validation is Recon-heavy; sensitivity check is reported.",
    ], left=1.0, top=4.0, width=10.8, height=2.0, size=19)
    add_speaker_note(slide, "Name the validation/test Recon gap before the committee does.")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Main Locked Result")
    add_table(slide, [
        ["Method", "Macro-F1", "Recon F1", "DE F1", "PR-AUC", "Rate"],
        ["Frozen MLP", "0.7685", "0.0250", "0.9157", "0.8732", "0.0000"],
        ["Locked selective TTA", "0.8658", "0.5050", "0.9202", "0.8738", "0.0470"],
    ], 0.75, 1.55, 11.8, 1.45, 15)
    add_bullets(slide, [
        "Macro-F1 delta: +0.0974.",
        "Recon F1 delta: +0.4800.",
        "PR-AUC delta: only +0.0006.",
        "Interpretation: operating-point decision-policy improvement.",
    ], top=3.45, size=21)
    add_speaker_note(slide, "Say clearly: this is not a broad ranking-quality gain.")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Alternative Explanation Check")
    add_bullets(slide, [
        "Critique: maybe TTA only filters uncertain rows.",
        "Control: reject the same fraction of frozen predictions.",
        "Matched reject rate: 4.7%.",
        "Result: frozen confidence rejection leaves Recon F1 at 0.0000.",
    ])
    add_table(slide, [
        ["Baseline", "Reject", "Macro", "Recon", "DE"],
        ["Frozen confidence reject", "0.0470", "0.7730", "0.0000", "0.9374"],
    ], 0.9, 5.3, 11.1, 0.9, 15)
    add_speaker_note(slide, "This is the strongest defense against the simplest critique.")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Defense Hardening")
    add_bullets(slide, [
        "Seven-seed robustness: Recon F1 0.5147 +/- 0.0589.",
        "Validation-distribution sensitivity mostly preserves the effect.",
        "Stronger frozen baselines still collapse on Recon.",
        "BN stream-order ablation checked.",
        "Override-to-Recon fraction: 0.8075; override-from-DE count: 0.",
    ], size=19)
    add_speaker_note(slide, "Hardening supports the locked story but does not replace the locked primary result.")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Operating-Point View")
    if RECON_FIG.exists():
        slide.shapes.add_picture(str(RECON_FIG), Inches(0.65), Inches(1.45), width=Inches(5.9))
    if DE_FIG.exists():
        slide.shapes.add_picture(str(DE_FIG), Inches(6.75), Inches(1.45), width=Inches(5.9))
    add_speaker_note(slide, "Use these figures to reinforce the decision-policy framing and DE safety discussion.")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Boundary Conditions")
    add_table(slide, [
        ["DAPT item", "Result"],
        ["MLP Macro F1", "0.6353 +/- 0.0043"],
        ["MLP Recon F1", "0.8932 +/- 0.0089"],
        ["TTA Macro delta", "-0.2874"],
        ["TTA Recon delta", "-0.6589"],
    ], 0.85, 1.45, 7.9, 2.6, 15)
    add_bullets(slide, [
        "DAPT supports detector-recipe transfer only.",
        "DAPT does not support TTA generality.",
        "DAPT DE support is too small for a DE safety claim.",
        "The boundary is visible by design.",
    ], left=9.0, top=1.55, width=3.7, height=3.7, size=17)
    add_speaker_note(slide, "This slide helps the committee trust the result because the boundary is visible.")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Final Defense Claim")
    add_textbox(slide, Inches(0.9), Inches(1.65), Inches(11.7), Inches(1.2), "Selective no-label TTA can recover rare-stage Reconnaissance under held-out source-file shift when constrained by a validation-selected safety gate.", 24, True)
    add_bullets(slide, [
        "Proves: locked rare-stage recovery and a safety-gated adaptation pattern.",
        "Does not prove: universal TTA, cross-dataset TTA, or adversarial-stream robustness.",
        "Next portfolio result: ATT&CK TTP-set profile retrieval.",
    ], top=3.4, size=21)
    add_speaker_note(slide, "End with discipline: narrow result, strong evidence, honest boundary.")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
